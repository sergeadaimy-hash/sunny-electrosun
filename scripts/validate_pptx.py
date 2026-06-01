"""Programmatic overlap + bounds validator for the case-study deck.

Reports:
- shapes that extend off-canvas (x+w > slide_w or y+h > slide_h or negative coords)
- non-background text boxes that overlap with each other by >5% of either's area
- backgrounds (full-slide rectangles) are excluded from overlap checks
"""
from pptx import Presentation
from pptx.util import Emu

PATH = "/Users/sergeadaimy/Desktop/Claude Projects/Sunny-Whatsapp Account Manager/presentation/sunny-case-study.pptx"
prs = Presentation(PATH)

slide_w = prs.slide_width
slide_h = prs.slide_height

def to_in(emu):
    return emu / 914400.0

def overlap_area(a, b):
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ox1 = max(ax1, bx1); oy1 = max(ay1, by1)
    ox2 = min(ax2, bx2); oy2 = min(ay2, by2)
    if ox2 <= ox1 or oy2 <= oy1:
        return 0
    return (ox2 - ox1) * (oy2 - oy1)

def shape_text(shp):
    try:
        if shp.has_text_frame:
            return (shp.text_frame.text or "").strip()[:60]
    except Exception:
        pass
    return ""

issues = 0
for idx, slide in enumerate(prs.slides, start=1):
    shapes = []
    for shp in slide.shapes:
        if shp.left is None or shp.top is None or shp.width is None or shp.height is None:
            continue
        x1, y1 = shp.left, shp.top
        x2, y2 = x1 + shp.width, y1 + shp.height
        text = shape_text(shp)
        # is_background: full-slide rectangle with a fill, no text, or near-full
        is_bg = (shp.width >= slide_w * 0.95 and shp.height >= slide_h * 0.95)
        shapes.append({
            "bbox": (x1, y1, x2, y2),
            "area": (x2 - x1) * (y2 - y1),
            "text": text,
            "is_bg": is_bg,
            "name": shp.name,
            "has_text": bool(text),
        })

    # off-canvas
    for sh in shapes:
        x1, y1, x2, y2 = sh["bbox"]
        if x1 < 0 or y1 < 0 or x2 > slide_w + 1000 or y2 > slide_h + 1000:
            print(f"slide {idx:02d} OFF-CANVAS: {sh['name']} at ({to_in(x1):.2f},{to_in(y1):.2f}) to ({to_in(x2):.2f},{to_in(y2):.2f})  text={sh['text']!r}")
            issues += 1

    # overlap: only between non-background text shapes
    text_shapes = [s for s in shapes if s["has_text"] and not s["is_bg"]]
    for i in range(len(text_shapes)):
        for j in range(i + 1, len(text_shapes)):
            a = text_shapes[i]
            b = text_shapes[j]
            ov = overlap_area(a["bbox"], b["bbox"])
            smaller = min(a["area"], b["area"]) or 1
            pct = ov / smaller
            if pct > 0.05:
                print(f"slide {idx:02d} OVERLAP {pct*100:.0f}%: {a['text']!r}  <>  {b['text']!r}")
                issues += 1

print(f"\nTOTAL ISSUES: {issues}")
