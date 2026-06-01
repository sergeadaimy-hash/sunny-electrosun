"""Build the Sunny case-study PPTX directly with python-pptx.

No HTML bridge. Every element is placed by hand with explicit non-overlapping
coordinates. Helvetica Neue throughout. Simple flat layouts. Real PPT shapes
for diagrams so PowerPoint stays fully editable.

Output: presentation/sunny-case-study.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Screenshot paths
SHOTS = Path("/Users/sergeadaimy/Downloads/Screenshots")
WA_PANELS_HOT  = str(SHOTS / "IMG_0545.PNG")    # solar panels, "I want to pay" → handoff
WA_PHOTO       = str(SHOTS / "IMG_0546.PNG")    # photo request: "I don't have a photo" → sends photo
WA_TECHNICAL   = str(SHOTS / "IMG_0547.PNG")    # 6kW inverter sizing + datasheet sent
WA_STOCK       = str(SHOTS / "IMG_0548.PNG")    # stock awareness: SE-F16 incoming, BOS-G out of stock
WA_BOS_HOT     = str(SHOTS / "IMG_0549.PNG")    # BOS-B 20 units + "yes i want to pay" → handoff

ADMIN = SHOTS / "Admin UI"
# macOS Screenshot tool uses U+202F (narrow no-break space) between time and AM/PM
_NBSP = " "
ADMIN_INBOX     = str(ADMIN / f"Screenshot 2026-05-27 at 7.04.50{_NBSP}PM.png")
ADMIN_CONTACTS  = str(ADMIN / f"Screenshot 2026-05-27 at 7.06.23{_NBSP}PM.png")
ADMIN_WAREHOUSE = str(ADMIN / f"Screenshot 2026-05-27 at 7.06.34{_NBSP}PM.png")
ADMIN_OWNER_CHAT_ONE  = str(ADMIN / f"Screenshot 2026-05-27 at 7.06.00{_NBSP}PM.png")
ADMIN_OWNER_CHAT_MANY = str(ADMIN / f"Screenshot 2026-05-27 at 7.06.07{_NBSP}PM.png")
ADMIN_RULES     = str(ADMIN / f"Screenshot 2026-05-27 at 7.06.56{_NBSP}PM.png")
ADMIN_MODELS    = str(ADMIN / f"Screenshot 2026-05-27 at 7.07.29{_NBSP}PM.png")

# -----------------------------------------------------------------------------
# Theme
# -----------------------------------------------------------------------------

CREAM       = RGBColor(0xf6, 0xef, 0xe1)
CREAM_DEEP  = RGBColor(0xef, 0xe5, 0xcf)
PAPER       = RGBColor(0xfb, 0xf6, 0xea)
INK         = RGBColor(0x15, 0x20, 0x2b)
INK_SOFT    = RGBColor(0x2d, 0x3a, 0x45)
INK_MUTE    = RGBColor(0x5a, 0x66, 0x72)
SUN         = RGBColor(0xf4, 0xa7, 0x2a)
SUN_DEEP    = RGBColor(0xe2, 0x87, 0x1d)
SUN_SOFT    = RGBColor(0xfd, 0xd9, 0x8a)
EMBER       = RGBColor(0xc5, 0x48, 0x1c)
TEAL        = RGBColor(0x1f, 0x6f, 0x7a)
TEAL_SOFT   = RGBColor(0xcd, 0xe5, 0xe6)
WHITE       = RGBColor(0xff, 0xff, 0xff)
LINE        = RGBColor(0xe5, 0xe0, 0xd0)
LINE_STRONG = RGBColor(0xcd, 0xc3, 0xa9)

HEAD_FONT = "Helvetica Neue"
BODY_FONT = "Helvetica Neue"

TOTAL = 42  # total slide count, used in page numbers


# -----------------------------------------------------------------------------
# Presentation setup
# -----------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout


# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------

def add_rect(slide, x, y, w, h, fill, line=None, line_w=0):
    """Add a plain rectangle. x/y/w/h in inches. fill: RGBColor or None for no fill."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w if line_w else 0.5)
    shp.shadow.inherit = False
    return shp


def add_round_rect(slide, x, y, w, h, fill, line=None, line_w=0.5, radius_in=0.12):
    """Add a rounded rectangle."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    # set corner radius via adjust
    try:
        shp.adjustments[0] = max(0.0, min(0.5, radius_in / min(w, h)))
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, text, x, y, w, h,
             size=18, color=INK, bold=False, italic=False,
             font=BODY_FONT, align=PP_ALIGN.LEFT, va=MSO_ANCHOR.TOP,
             line_spacing=1.15, letter_spacing_pts=None):
    """Add a text box. All sizes in inches, font size in pt."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = va

    # Support multi-line text
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if letter_spacing_pts is not None:
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(int(letter_spacing_pts * 100)))
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, weight=0.75):
    """Add a straight line. coords in inches, weight in pt."""
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_arrow(slide, x1, y1, x2, y2, color=INK, weight=1.25):
    """Add a connector with an arrowhead at the end."""
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    # add arrow end
    lnElem = ln.line._get_or_add_ln()
    tail = etree.SubElement(lnElem, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return ln


def add_pill(slide, x, y, w, h, text, fill=SUN_SOFT, text_color=INK, size=11, bold=False):
    """Pill-shaped badge."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = BODY_FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    return shp


def new_slide(page_num, show_chrome=True, bg=CREAM):
    """Create a blank slide with cream background and optional bottom chrome."""
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.333, 7.5, bg, line=None)
    if show_chrome:
        add_text(slide, "SUNNY  ·  CASE STUDY  ·  ELECTROSUN", 0.7, 7.10, 7, 0.25,
                 size=9, color=INK_MUTE, font=BODY_FONT, align=PP_ALIGN.LEFT, bold=True)
        add_text(slide, f"{page_num:02d} / {TOTAL:02d}", 11.633, 7.10, 1.0, 0.25,
                 size=9, color=INK_MUTE, font=BODY_FONT, align=PP_ALIGN.RIGHT)
        # subtle hairline above chrome
        add_line(slide, 0.7, 7.05, 12.633, 7.05, color=LINE, weight=0.5)
    return slide


def add_eyebrow(slide, label):
    """Eyebrow tag at top of content slide."""
    # ember dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(0.52), Inches(0.16), Inches(0.16))
    dot.fill.solid()
    dot.fill.fore_color.rgb = EMBER
    dot.line.fill.background()
    dot.shadow.inherit = False
    add_text(slide, label.upper(), 0.95, 0.52, 8, 0.22,
             size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)


def add_headline(slide, text, size=40, y=1.0, x=0.7, w=11.933, h=1.6, color=INK):
    add_text(slide, text, x, y, w, h, size=size, color=color, bold=True, font=HEAD_FONT,
             line_spacing=1.08)


def add_body(slide, text, x=0.7, y=2.7, w=11.933, h=4.0, size=18, color=INK_SOFT):
    add_text(slide, text, x, y, w, h, size=size, color=color, line_spacing=1.35)


# -----------------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------------

def slide_01_cover():
    s = new_slide(1, show_chrome=False)
    # Cream + subtle sun band on right
    add_rect(s, 9.5, 0, 3.833, 7.5, CREAM_DEEP, line=None)
    # small sun icon top right
    sun = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(0.6), Inches(0.55), Inches(0.55))
    sun.fill.solid(); sun.fill.fore_color.rgb = SUN; sun.line.fill.background(); sun.shadow.inherit = False
    # Pre-title
    add_text(s, "CASE STUDY", 0.7, 1.4, 6, 0.3, size=12, color=EMBER, bold=True, letter_spacing_pts=4.0)
    # Big name (Helvetica Neue bold italic, sized to fit without descender bleed)
    add_text(s, "Sunny.", 0.7, 1.9, 9, 2.0, size=110, color=INK, bold=True, font=HEAD_FONT, italic=True, line_spacing=1.0)
    # Subtitle, generous space below to clear descenders
    add_text(s, "An AI Account Manager",
             0.7, 4.25, 11, 0.6, size=34, color=INK, font=HEAD_FONT, bold=True, line_spacing=1.1)
    add_text(s, "for WhatsApp Business.",
             0.7, 4.95, 11, 0.6, size=34, color=INK_SOFT, font=HEAD_FONT, line_spacing=1.1)
    # Client line
    add_text(s, "Built for ElectroSun, Nigeria.",
             0.7, 5.7, 11, 0.45, size=18, color=INK_SOFT)
    # Footer
    add_line(s, 0.7, 6.6, 12.633, 6.6, color=LINE_STRONG, weight=0.5)
    add_text(s, "Built by [Your Agency]", 0.7, 6.75, 6, 0.3, size=12, color=INK_MUTE, bold=True)
    add_text(s, "May 2026", 6.7, 6.75, 5.933, 0.3, size=12, color=INK_MUTE, align=PP_ALIGN.RIGHT)


def slide_02_pitch():
    s = new_slide(2)
    add_eyebrow(s, "01  ·  the pitch")
    add_headline(s, "Sunny replies to every WhatsApp customer in seconds,\nin five languages, and knows when to call the boss.",
                 size=34, h=2.0)
    # three stat tiles
    tiles = [
        ("8 sec", "average reply time", "Across text, voice notes, and photos."),
        ("5", "languages spoken", "English, Pidgin, Hausa, Yoruba, Igbo."),
        ("100%", "of leads classified", "HOT, SERIOUS, COLD, REPEAT, DISQUALIFIED."),
    ]
    for i, (big, label, body) in enumerate(tiles):
        x = 0.7 + i * 4.1
        add_round_rect(s, x, 3.8, 3.85, 2.6, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
        add_text(s, big, x + 0.3, 4.0, 3.4, 1.1, size=54, color=INK, bold=True, font=HEAD_FONT, line_spacing=1.0)
        add_text(s, label, x + 0.3, 5.1, 3.4, 0.35, size=13, color=EMBER, bold=True, letter_spacing_pts=1.5)
        add_text(s, body, x + 0.3, 5.5, 3.4, 0.85, size=14, color=INK_SOFT, line_spacing=1.35)
    # footer line
    add_text(s, "He doesn't sleep. He doesn't forget. He never invents a price.",
             0.7, 6.75, 11.933, 0.3, size=14, color=INK_MUTE, italic=True)


def slide_03_opportunity():
    s = new_slide(3)
    add_eyebrow(s, "02  ·  the opportunity")
    add_headline(s, "WhatsApp is where Africa shops.\nMost small businesses can't keep up.", size=36, h=2.0)
    # Two columns of body
    add_text(s, "2.7 BILLION USERS WORLDWIDE",
             0.7, 3.5, 5.7, 0.3, size=11, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s,
             "In Nigeria, ninety percent of business inquiries arrive on WhatsApp first. "
             "Most are answered slowly. Many are answered never.",
             0.7, 3.85, 5.7, 1.8, size=18, color=INK_SOFT)
    add_text(s, "ELECTROSUN'S REALITY",
             6.93, 3.5, 5.7, 0.3, size=11, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s,
             "Leads piled up faster than the team could reply. Customers left for competitors. "
             "ElectroSun needed a 24/7 account manager who knew the catalog and spoke five languages.",
             6.93, 3.85, 5.7, 2.2, size=18, color=INK_SOFT)


def slide_04_before():
    s = new_slide(4)
    add_eyebrow(s, "03  ·  before sunny")
    add_headline(s, "Before Sunny, leads piled up faster than the team\ncould reply.", size=34, h=1.8)
    # Timeline-style row of events
    events = [
        ("09:00", "Customer asks for inverter quote"),
        ("11:00", "Message still unread"),
        ("14:00", "Staff finally replies"),
        ("14:15", "Customer already bought elsewhere"),
    ]
    y = 3.6
    for i, (time, text) in enumerate(events):
        x = 0.7 + i * 3.1
        # dot on timeline
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.5), Inches(y + 0.15), Inches(0.22), Inches(0.22))
        dot.fill.solid(); dot.fill.fore_color.rgb = (EMBER if i == 3 else SUN); dot.line.fill.background(); dot.shadow.inherit = False
        # time
        add_text(s, time, x, y + 0.5, 2.8, 0.4, size=22, color=INK, bold=True, font=HEAD_FONT)
        # text
        add_text(s, text, x, y + 1.0, 2.8, 1.2, size=14, color=INK_SOFT, line_spacing=1.3)
    # connecting line
    add_line(s, 1.05, y + 0.26, 11.0, y + 0.26, color=LINE_STRONG, weight=1.0)
    # bottom stat row
    add_text(s, "AVG FIRST REPLY: 3+ HOURS  ·  VOICE NOTES IGNORED: MOST  ·  LOST LEADS PER WEEK: DOZENS",
             0.7, 6.5, 11.933, 0.3, size=11, color=INK_MUTE, bold=True, letter_spacing_pts=2.0)


def slide_05_chatbot_fails():
    s = new_slide(5)
    add_eyebrow(s, "04  ·  the chatbot trap")
    add_headline(s, "Off-the-shelf chatbots fail at exactly\nthe moments that matter.", size=34, h=1.8)
    fails = [
        ("01", "No memory",        "Forgets what the customer said two messages ago."),
        ("02", "No business context", "Quotes generic prices, can't see the warehouse."),
        ("03", "No escalation",    "Has no idea when the boss needs to be called."),
        ("04", "One language",     "Misses Pidgin, Hausa, Yoruba, Igbo customers."),
        ("05", "No learning",      "Same mistakes, every conversation."),
    ]
    y = 3.5
    for i, (num, title, body) in enumerate(fails):
        rowy = y + i * 0.66
        add_text(s, num, 0.7, rowy + 0.08, 0.7, 0.4, size=22, color=SUN_DEEP, bold=True, font=HEAD_FONT)
        add_text(s, title, 1.7, rowy + 0.05, 3.5, 0.5, size=20, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body, 5.5, rowy + 0.08, 7.0, 0.5, size=16, color=INK_SOFT)


def slide_06_capabilities():
    s = new_slide(6)
    add_eyebrow(s, "05  ·  capabilities")
    add_headline(s, "Eight things Sunny does, autonomously.", size=36)
    caps = [
        ("Reply in seconds",       "Text, voice, photos."),
        ("Quote from live stock",  "Never invents a price."),
        ("Classify every lead",    "HOT  ·  SERIOUS  ·  COLD."),
        ("Escalate to the owner",  "Payment-ready, negotiations, complaints."),
        ("Transcribe voice notes", "Whisper-powered, every language."),
        ("Send photos & datasheets","On request, from warehouse."),
        ("Speak 5 languages",      "English, Pidgin, Hausa, Yoruba, Igbo."),
        ("Report to the boss",     "Daily summaries, on-demand Q&A."),
    ]
    cols, rows = 4, 2
    cell_w, cell_h = 2.95, 1.65
    start_x, start_y = 0.7, 3.4
    gap = 0.15
    for i, (title, body) in enumerate(caps):
        r, c = divmod(i, cols)
        x = start_x + c * (cell_w + gap)
        y = start_y + r * (cell_h + gap)
        add_round_rect(s, x, y, cell_w, cell_h, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.12)
        # small sun dot
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(y + 0.27), Inches(0.18), Inches(0.18))
        d.fill.solid(); d.fill.fore_color.rgb = SUN; d.line.fill.background(); d.shadow.inherit = False
        add_text(s, title, x + 0.25, y + 0.55, cell_w - 0.5, 0.42, size=16, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body,  x + 0.25, y + 1.05, cell_w - 0.5, 0.55, size=12, color=INK_SOFT, line_spacing=1.3)


def slide_07_not_chatbot():
    s = new_slide(7)
    add_eyebrow(s, "06  ·  the difference")
    add_headline(s, "Sunny isn't a chatbot. He's an account manager.", size=34)
    # 5 rows: dimension, chatbot, sunny
    rows = [
        ("Memory",     "Stateless",                 "Tracks every fact the customer has shared"),
        ("Knowledge",  "Generic FAQ",               "Live warehouse, prices, datasheets, photos"),
        ("Escalation", "None",                      "Pages the owner with full context"),
        ("Languages",  "One",                       "Five, including Nigerian Pidgin"),
        ("Self-check", "None",                      "14 reply guards before every send"),
    ]
    y0 = 3.0
    # header
    add_text(s, "DIMENSION", 0.7, y0, 2.5, 0.3, size=10, color=INK_MUTE, bold=True, letter_spacing_pts=2.0)
    add_text(s, "CHATBOT",   3.7, y0, 4.0, 0.3, size=10, color=INK_MUTE, bold=True, letter_spacing_pts=2.0)
    add_text(s, "SUNNY",     8.0, y0, 4.6, 0.3, size=10, color=EMBER,    bold=True, letter_spacing_pts=2.0)
    add_line(s, 0.7, y0 + 0.35, 12.633, y0 + 0.35, color=LINE_STRONG, weight=0.75)
    for i, (dim, cb, sn) in enumerate(rows):
        ry = y0 + 0.55 + i * 0.65
        add_text(s, dim, 0.7, ry, 2.5, 0.5, size=16, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, cb,  3.7, ry, 4.0, 0.5, size=15, color=INK_MUTE)
        add_text(s, sn,  8.0, ry, 4.6, 0.5, size=15, color=INK)
        if i < len(rows) - 1:
            add_line(s, 0.7, ry + 0.55, 12.633, ry + 0.55, color=LINE, weight=0.4)


def slide_08_day():
    s = new_slide(8)
    add_eyebrow(s, "07  ·  a day in his life")
    add_headline(s, "A typical Tuesday.", size=40)
    events = [
        ("06:14", "Hausa voice note about a 5 kVA inverter.", "Sunny transcribes, finds the model, replies in Hausa."),
        ("09:02", "Customer asks for a product brochure.", "Sunny matches the item and sends the datasheet PDF."),
        ("11:47", "Customer pushes: 'best price?'", "Sunny answers honestly: 'Yes, this is our best price.'"),
        ("12:30", "Customer: 'send me your account, I want to pay.'", "Sunny classifies HOT and pages the owner with context."),
        ("15:21", "Customer sends a photo: 'what is this?'", "Sunny uses vision, identifies the brand, replies."),
        ("18:00", "Owner asks Sunny: 'how many hot leads today?'", "Sunny answers from the live database."),
    ]
    y0 = 2.9
    # vertical timeline line
    add_line(s, 1.8, y0 - 0.05, 1.8, y0 + 4.2, color=LINE_STRONG, weight=1.0)
    for i, (time, what, action) in enumerate(events):
        ry = y0 + i * 0.7
        # dot
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.62), Inches(ry + 0.05), Inches(0.36), Inches(0.36))
        d.fill.solid(); d.fill.fore_color.rgb = SUN; d.line.color.rgb = SUN_DEEP; d.line.width = Pt(0.75); d.shadow.inherit = False
        add_text(s, time, 0.7, ry + 0.08, 1.0, 0.4, size=14, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, what, 2.3, ry + 0.05, 5.5, 0.45, size=15, color=INK, bold=True)
        add_text(s, action, 7.95, ry + 0.05, 4.7, 0.45, size=14, color=INK_SOFT)


def slide_09_numbers():
    s = new_slide(9)
    add_eyebrow(s, "08  ·  the numbers")
    add_headline(s, "What Sunny does, at scale.", size=40)
    stats = [
        ("6 sec",   "Average text reply time"),
        ("42 sec",  "Average voice-note reply time"),
        ("5",       "Languages spoken"),
        ("24 / 7",  "Always on, no shifts"),
        ("$0.02",   "Avg cost per customer message"),
        ("0",       "Hallucinated prices (14 guards)"),
    ]
    cols, rows = 3, 2
    cell_w, cell_h = 3.95, 1.8
    start_x, start_y = 0.7, 3.0
    gap_x, gap_y = 0.18, 0.2
    for i, (big, label) in enumerate(stats):
        r, c = divmod(i, cols)
        x = start_x + c * (cell_w + gap_x)
        y = start_y + r * (cell_h + gap_y)
        add_round_rect(s, x, y, cell_w, cell_h, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
        add_text(s, big, x + 0.3, y + 0.25, cell_w - 0.6, 1.0, size=58, color=INK, bold=True, font=HEAD_FONT, line_spacing=1.0)
        add_text(s, label, x + 0.3, y + 1.3, cell_w - 0.6, 0.4, size=14, color=INK_SOFT)


def slide_10_architecture():
    s = new_slide(10)
    add_eyebrow(s, "09  ·  architecture")
    add_headline(s, "How Sunny is wired.", size=40)

    # Three columns: WhatsApp Cloud API | Railway container | External services
    # Left: WhatsApp
    add_round_rect(s, 0.7, 3.3, 2.9, 3.0, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.12)
    add_text(s, "WHATSAPP", 0.85, 3.45, 2.6, 0.25, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s, "Cloud API", 0.85, 3.78, 2.6, 0.5, size=22, color=INK, bold=True, font=HEAD_FONT)
    add_text(s, "Official Meta channel.\nSigned webhook in.\nText, voice, photos, calls.",
             0.85, 4.4, 2.6, 1.7, size=12, color=INK_SOFT, line_spacing=1.4)

    # Middle: Railway container
    add_round_rect(s, 4.1, 2.7, 5.1, 4.0, WHITE, line=LINE_STRONG, line_w=1.0, radius_in=0.15)
    add_text(s, "RAILWAY CONTAINER", 4.25, 2.85, 4.85, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    pipeline = [
        ("Webhook + signature verify",  EMBER),
        ("Debounce queue (6 sec)",      SUN_DEEP),
        ("Classifier (Sonnet 4.6)",     SUN),
        ("Reply engine (Opus 4.7)",     SUN),
        ("14 reply guards",             TEAL),
        ("Send + persist",              INK_SOFT),
    ]
    y = 3.2
    for label, c in pipeline:
        add_round_rect(s, 4.3, y, 4.75, 0.36, CREAM, line=LINE_STRONG, line_w=0.3, radius_in=0.08)
        # color tab
        tab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.3), Inches(y), Inches(0.08), Inches(0.36))
        tab.fill.solid(); tab.fill.fore_color.rgb = c; tab.line.fill.background(); tab.shadow.inherit = False
        add_text(s, label, 4.5, y + 0.05, 4.5, 0.27, size=12, color=INK, va=MSO_ANCHOR.MIDDLE)
        y += 0.43
    # storage row
    add_round_rect(s, 4.3, 5.85, 2.3, 0.6, CREAM_DEEP, line=LINE_STRONG, line_w=0.3, radius_in=0.08)
    add_text(s, "SQLite DB", 4.45, 5.93, 2.0, 0.4, size=12, color=INK, bold=True, va=MSO_ANCHOR.MIDDLE)
    add_round_rect(s, 6.75, 5.85, 2.3, 0.6, CREAM_DEEP, line=LINE_STRONG, line_w=0.3, radius_in=0.08)
    add_text(s, "Media volume", 6.9, 5.93, 2.0, 0.4, size=12, color=INK, bold=True, va=MSO_ANCHOR.MIDDLE)

    # Right: External services
    add_round_rect(s, 9.75, 3.3, 2.9, 3.0, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.12)
    add_text(s, "EXTERNAL APIs", 9.9, 3.45, 2.6, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s, "Claude API (Anthropic)", 9.9, 3.85, 2.6, 0.35, size=14, color=INK, bold=True)
    add_text(s, "Whisper API (OpenAI)",   9.9, 4.25, 2.6, 0.35, size=14, color=INK, bold=True)
    add_text(s, "Owner WhatsApp",          9.9, 4.65, 2.6, 0.35, size=14, color=INK, bold=True)
    add_text(s, "GitHub (code source)",    9.9, 5.05, 2.6, 0.35, size=14, color=INK, bold=True)

    # arrows
    add_arrow(s, 3.6, 4.7, 4.05, 4.7, color=INK_SOFT, weight=1.5)
    add_arrow(s, 9.2, 4.5, 9.7, 4.3, color=INK_SOFT, weight=1.5)
    add_arrow(s, 9.2, 4.9, 9.7, 5.1, color=INK_SOFT, weight=1.5)

    add_text(s, "GitHub holds the code. Railway runs the agent and stores the data. WhatsApp Cloud API is the channel.",
             0.7, 6.75, 11.933, 0.3, size=12, color=INK_MUTE, italic=True)


def slide_11_where_lives():
    s = new_slide(11)
    add_eyebrow(s, "10  ·  where he lives")
    add_headline(s, "Three places. Three jobs.", size=40)
    cols = [
        ("GitHub",       "SOURCE OF TRUTH FOR CODE",
         "Every change is a commit. Every commit is auditable. If a laptop dies, nothing is lost. The system rebuilds in minutes."),
        ("Railway",      "THE RUNNING AGENT",
         "Container plus SQLite plus media volume. Auto-redeploys on every push to main. This is where Sunny actually thinks and speaks."),
        ("WhatsApp",     "THE CHANNEL",
         "Official Meta Cloud API. Verified business number. Signed webhooks. No third-party wrappers."),
    ]
    for i, (title, sub, body) in enumerate(cols):
        x = 0.7 + i * 4.1
        add_round_rect(s, x, 3.0, 3.85, 3.6, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
        # accent number
        add_text(s, f"0{i+1}", x + 0.3, 3.2, 1.0, 0.3, size=14, color=EMBER, bold=True, letter_spacing_pts=2.0)
        add_text(s, title, x + 0.3, 3.6, 3.4, 0.6, size=32, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, sub,   x + 0.3, 4.3, 3.4, 0.3, size=10, color=INK_MUTE, bold=True, letter_spacing_pts=1.5)
        add_text(s, body,  x + 0.3, 4.7, 3.4, 1.9, size=14, color=INK_SOFT, line_spacing=1.4)


def slide_12_brain():
    s = new_slide(12)
    add_eyebrow(s, "11  ·  the brain")
    add_headline(s, "One model across the whole pipeline.", size=36, h=1.4)
    # Big claim
    add_round_rect(s, 0.7, 3.0, 11.933, 1.4, INK, line=None, radius_in=0.18)
    add_text(s, "Claude Sonnet 4.6", 0.95, 3.15, 11.5, 0.7, size=44, color=SUN_SOFT, bold=True, italic=True, font=HEAD_FONT, va=MSO_ANCHOR.MIDDLE)
    add_text(s, "Anthropic  ·  cached prompts  ·  vision  ·  long context",
             0.95, 3.85, 11.5, 0.45, size=14, color=CREAM, va=MSO_ANCHOR.MIDDLE, letter_spacing_pts=1.0)
    # Three call sites
    sites = [
        ("Reply",      "Customer-facing messages.\nVoice rules, 14 guards."),
        ("Classifier", "Every inbound is categorized.\nHOT, SERIOUS, COLD."),
        ("Owner Q&A",  "Boss-mode questions.\nLive data snapshot."),
    ]
    for i, (name, body) in enumerate(sites):
        x = 0.7 + i * 4.1
        add_round_rect(s, x, 4.65, 3.85, 1.75, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.12)
        add_text(s, name, x + 0.3, 4.8, 3.4, 0.5, size=20, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body, x + 0.3, 5.3, 3.4, 1.05, size=13, color=INK_SOFT, line_spacing=1.4)
    add_text(s, "Sonnet is the sweet spot: rule discipline of a top model at roughly one-fifth the cost of frontier.",
             0.7, 6.6, 11.933, 0.3, size=12, color=INK_MUTE, italic=True, align=PP_ALIGN.CENTER)


def slide_13_pipeline():
    s = new_slide(13)
    add_eyebrow(s, "12  ·  the pipeline")
    add_headline(s, "What happens when a customer sends a message.", size=34)
    steps = [
        "Inbound webhook",
        "Idempotency check",
        "Media handling",
        "Debounce queue",
        "Classify (Sonnet)",
        "Reply (Opus)",
        "14 guards",
        "Send + persist",
    ]
    box_w, box_h = 1.35, 1.0
    gap = 0.13
    start_x = 0.7
    y = 3.7
    total_w = len(steps) * box_w + (len(steps) - 1) * gap
    # center the row
    start_x = (13.333 - total_w) / 2
    for i, label in enumerate(steps):
        x = start_x + i * (box_w + gap)
        add_round_rect(s, x, y, box_w, box_h, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.1)
        add_text(s, f"{i+1:02d}", x + 0.1, y + 0.1, box_w - 0.2, 0.25, size=10, color=EMBER, bold=True, letter_spacing_pts=1.5)
        add_text(s, label, x + 0.1, y + 0.42, box_w - 0.2, 0.6, size=11, color=INK, bold=True, va=MSO_ANCHOR.TOP, line_spacing=1.2)
        # arrow to next
        if i < len(steps) - 1:
            ax1 = x + box_w + 0.01
            ax2 = x + box_w + gap - 0.01
            add_arrow(s, ax1, y + box_h / 2, ax2, y + box_h / 2, color=INK_SOFT, weight=1.5)
    add_text(s, "Idempotent at every step. A retried webhook will never double-send.",
             0.7, 5.6, 11.933, 0.4, size=14, color=INK_SOFT, italic=True, align=PP_ALIGN.CENTER)


def slide_14_state():
    s = new_slide(14)
    add_eyebrow(s, "13  ·  memory")
    add_headline(s, "Sunny doesn't ask the same question twice.", size=36)
    # left explanation
    add_text(s, "Before every reply, Sunny builds a structured state object that captures what the customer has shared, what's already been asked, and what's still pending. The object is injected into the prompt as a system block.",
             0.7, 3.0, 5.6, 3.5, size=16, color=INK_SOFT, line_spacing=1.4)
    # right: stylized "code block"
    add_round_rect(s, 6.8, 3.0, 5.83, 3.5, INK, line=None, radius_in=0.15)
    # tab dots
    for j, c in enumerate([RGBColor(0xff,0x5f,0x56), RGBColor(0xff,0xbd,0x2e), RGBColor(0x27,0xc9,0x3f)]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.95 + j*0.22), Inches(3.18), Inches(0.13), Inches(0.13))
        d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background(); d.shadow.inherit = False
    code = (
        "conversation_state = {\n"
        "  facts: {\n"
        "    kw: 16,\n"
        "    phase: '3-phase',\n"
        "    location: 'Lagos',\n"
        "    use_case: 'home'\n"
        "  },\n"
        "  asked_already: [\n"
        "    'installer_or_end_user',\n"
        "    'budget_range'\n"
        "  ],\n"
        "  in_current_msg: [\n"
        "    'battery_price'\n"
        "  ]\n"
        "}"
    )
    add_text(s, code, 6.95, 3.55, 5.6, 2.85, size=12, color=SUN_SOFT,
             font="Menlo", line_spacing=1.35)
    add_text(s, "Result: no 'and how big is your system?' five times. Conversation moves forward, every turn.",
             0.7, 6.75, 11.933, 0.3, size=12, color=INK_MUTE, italic=True)


def slide_15_guards():
    s = new_slide(15)
    add_eyebrow(s, "14  ·  self-check")
    add_headline(s, "Every reply runs through 14 validators before it ships.", size=30, h=1.4)
    guards = [
        "Price-dump guard",       "Repeat guard",
        "Trailing-question strip","CTA-tail strip",
        "wa.me URL strip",        "Prompt-leak detector",
        "Owner-number leak",      "Phone-list-dump block",
        "Catalog enumeration",    "SKU list dump block",
        "Fabricated variant",     "HV BOM validator",
        "BOM cleanup pass",       "No-double-dashes",
    ]
    cols, rows_c = 2, 7
    cell_w, cell_h = 5.95, 0.42
    start_x, start_y = 0.7, 2.9
    gap_x, gap_y = 0.15, 0.12
    for i, g in enumerate(guards):
        r, c = divmod(i, cols)
        x = start_x + c * (cell_w + gap_x)
        y = start_y + r * (cell_h + gap_y)
        add_round_rect(s, x, y, cell_w, cell_h, PAPER, line=LINE_STRONG, line_w=0.4, radius_in=0.08)
        # number
        add_text(s, f"{i+1:02d}", x + 0.18, y + 0.05, 0.5, 0.32, size=12, color=EMBER, bold=True, va=MSO_ANCHOR.MIDDLE)
        add_text(s, g, x + 0.7, y + 0.05, cell_w - 0.85, 0.32, size=14, color=INK, va=MSO_ANCHOR.MIDDLE)
    add_text(s, "Nothing reaches the customer until all 14 pass.",
             0.7, 6.75, 11.933, 0.3, size=12, color=INK_MUTE, italic=True, align=PP_ALIGN.CENTER)


def slide_16_warehouse():
    s = new_slide(16)
    add_eyebrow(s, "15  ·  source of truth")
    add_headline(s, "The warehouse is the only place Sunny gets prices.", size=32, h=1.4)
    # Left: explanation
    add_text(s, "Owner edits live from the admin panel. Each item has brand, model, price, stock state per location, optional datasheet PDF, and optional product photos.",
             0.7, 2.9, 5.6, 2.0, size=16, color=INK_SOFT, line_spacing=1.4)
    add_text(s, "Changes hit the next reply. No redeploy.",
             0.7, 5.0, 5.6, 0.5, size=16, color=EMBER, italic=True, bold=True)
    # Right: mock warehouse card
    add_round_rect(s, 6.95, 2.9, 5.7, 3.85, WHITE, line=LINE_STRONG, line_w=0.75, radius_in=0.15)
    add_text(s, "WAREHOUSE ITEM", 7.1, 3.05, 5.4, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=1.5)
    add_text(s, "Sungrow 12kW 3-phase Inverter", 7.1, 3.4, 5.4, 0.6, size=22, color=INK, bold=True, font=HEAD_FONT)
    # location rows
    add_text(s, "ABUJA",  7.1, 4.15, 1.5, 0.22, size=10, color=INK_MUTE, bold=True, letter_spacing_pts=1.5)
    add_text(s, "In stock  ·  4 units", 7.1, 4.45, 5.4, 0.35, size=16, color=INK)
    add_text(s, "LAGOS",  7.1, 4.92, 1.5, 0.22, size=10, color=INK_MUTE, bold=True, letter_spacing_pts=1.5)
    add_text(s, "Incoming  ·  ETA Jun 10", 7.1, 5.18, 5.4, 0.35, size=16, color=INK)
    # price + attachments
    add_text(s, "PRICE", 7.1, 5.65, 1.5, 0.22, size=10, color=INK_MUTE, bold=True, letter_spacing_pts=1.5)
    add_text(s, "₦4,850,000", 7.1, 5.92, 5.4, 0.4, size=22, color=EMBER, bold=True, font=HEAD_FONT)
    add_pill(s, 7.1, 6.4, 1.5, 0.3, "Datasheet", fill=SUN_SOFT, size=10, bold=True)
    add_pill(s, 8.7, 6.4, 1.5, 0.3, "3 photos", fill=TEAL_SOFT, size=10, bold=True)


def slide_17_knowledge():
    s = new_slide(17)
    add_eyebrow(s, "16  ·  knowledge")
    add_headline(s, "Four layers feed every reply.", size=40)
    layers = [
        ("01", "Master prompt",        "Personality, voice rules, 19 sections of doctrine."),
        ("02", "Warehouse stock",      "Live inventory, prices, datasheets, photos."),
        ("03", "Conversation history", "Last 50 messages, alternating roles, deduped."),
        ("04", "Conversation state",   "Structured facts, asked-already, current questions."),
    ]
    y = 3.3
    for i, (num, title, body) in enumerate(layers):
        ry = y + i * 0.78
        add_round_rect(s, 0.7, ry, 11.933, 0.7, PAPER, line=LINE_STRONG, line_w=0.4, radius_in=0.1)
        add_text(s, num, 0.95, ry + 0.15, 0.7, 0.4, size=18, color=EMBER, bold=True, font=HEAD_FONT)
        add_text(s, title, 1.85, ry + 0.12, 3.8, 0.45, size=18, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body, 5.95, ry + 0.15, 6.6, 0.4, size=14, color=INK_SOFT, va=MSO_ANCHOR.MIDDLE)
    add_text(s, "Every reply call carries all four. Prompt caching keeps cost down.",
             0.7, 6.75, 11.933, 0.3, size=12, color=INK_MUTE, italic=True)


def slide_18_owner_qa():
    s = new_slide(18)
    add_eyebrow(s, "17  ·  boss mode")
    add_headline(s, "The owner can ask Sunny anything about his own business.", size=28, h=1.2)
    # Mock placeholder
    add_round_rect(s, 0.7, 2.7, 11.933, 4.0, WHITE, line=LINE_STRONG, line_w=0.75, radius_in=0.15)
    # [MOCK] badge
    add_round_rect(s, 11.95, 2.78, 0.6, 0.28, EMBER, line=None, radius_in=0.14)
    add_text(s, "MOCK", 11.95, 2.78, 0.6, 0.28, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE, letter_spacing_pts=1.5)

    add_text(s, "OWNER  →  SUNNY", 0.95, 2.9, 6, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s, "How many hot leads today?", 0.95, 3.2, 11.5, 0.55, size=20, color=INK, bold=True, font=HEAD_FONT)
    add_text(s, "SUNNY  →  OWNER", 0.95, 3.95, 6, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s,
             "Three hot leads today.\n\n"
             "1.  Adebayo  ·  16kW residential, asked for the account at 11:47.\n"
             "2.  Chinedu  ·  3-phase commercial, wants a site visit, last message 14:02.\n"
             "3.  Funke   ·  battery only, ready to pay, last message 17:30.",
             0.95, 4.25, 11.5, 2.25, size=15, color=INK_SOFT, line_spacing=1.45)


def slide_19_multimodal():
    s = new_slide(19)
    add_eyebrow(s, "18  ·  multimodal")
    add_headline(s, "Voice notes, photos, datasheets, all handled.", size=34)
    cols = [
        ("Voice",      "OpenAI Whisper transcribes every voice note before classification. The customer speaks Pidgin or Hausa; Sunny understands and replies in the same language."),
        ("Photos",     "Customer sends a product photo: Sunny uses vision to identify the brand. Customer asks for a photo: Sunny sends up to 3 images from warehouse."),
        ("Documents",  "Customer asks for a brochure: Sunny matches the warehouse item and sends the PDF. Cached on Meta's CDN for 25 days to save reuploads."),
    ]
    for i, (title, body) in enumerate(cols):
        x = 0.7 + i * 4.1
        add_round_rect(s, x, 3.1, 3.85, 3.7, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
        # icon dot
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.35), Inches(3.4), Inches(0.4), Inches(0.4))
        d.fill.solid(); d.fill.fore_color.rgb = SUN; d.line.fill.background(); d.shadow.inherit = False
        add_text(s, title, x + 0.95, 3.4, 2.9, 0.5, size=26, color=INK, bold=True, font=HEAD_FONT, va=MSO_ANCHOR.MIDDLE)
        add_text(s, body, x + 0.35, 4.1, 3.4, 2.5, size=14, color=INK_SOFT, line_spacing=1.45)


def slide_20_lead_cats():
    s = new_slide(20)
    add_eyebrow(s, "19  ·  lead categories")
    add_headline(s, "Every conversation lands in one of five buckets.", size=32)
    cats = [
        ("HOT",            EMBER,  "Asked for the account.\nReady to pay."),
        ("SERIOUS",        SUN_DEEP,"Specific product.\nReal budget, real timeline."),
        ("COLD",           INK_MUTE,"Browsing.\nGeneral questions."),
        ("REPEAT_CLIENT",  TEAL,   "Existing customer\nasking again."),
        ("DISQUALIFIED",   INK_SOFT,"Not a real prospect."),
    ]
    cell_w, cell_h = 2.32, 3.0
    gap = 0.16
    start_x = 0.7
    y = 3.1
    for i, (name, color, body) in enumerate(cats):
        x = start_x + i * (cell_w + gap)
        add_round_rect(s, x, y, cell_w, cell_h, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
        # color band on top
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cell_w), Inches(0.35))
        band.fill.solid(); band.fill.fore_color.rgb = color; band.line.fill.background(); band.shadow.inherit = False
        add_text(s, name, x + 0.2, y + 0.55, cell_w - 0.4, 0.5, size=18, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body, x + 0.2, y + 1.2, cell_w - 0.4, 1.7, size=13, color=INK_SOFT, line_spacing=1.45)
    add_text(s, "Owner sees the category live in admin. Reports group by category.",
             0.7, 6.75, 11.933, 0.3, size=12, color=INK_MUTE, italic=True)


def slide_21_hot_handoff():
    s = new_slide(21)
    add_eyebrow(s, "20  ·  hot lead")
    add_headline(s, "When the customer says 'I want to pay', the owner is paged in seconds.", size=22, h=1.0)
    # Customer side (left): real WhatsApp screenshot
    add_image_card(s, WA_BOS_HOT, 0.7, 2.35, 3.6, 4.05, label="CUSTOMER  ·  WHATSAPP")
    # Owner side (right): real admin Owner Chat
    add_image_card(s, ADMIN_OWNER_CHAT_ONE, 4.5, 2.35, 8.13, 4.05, label="OWNER  ·  ADMIN  ·  HOT ALERT")
    # explanation footer
    add_text(s,
             "Left: customer confirms 20 BOS-B units across both warehouses, then says 'yes i want to pay'. "
             "Sunny replies with one short line and the specialist's direct WhatsApp link. He never invents account numbers or invoice figures.   "
             "Right: within seconds, a HOT LEAD alert lands in Owner Chat with customer name, phone, lead category, intent (payment_question), the full conversation summary, and a click-to-chat link.",
             0.7, 6.55, 11.933, 0.5, size=10, color=INK_MUTE, italic=True, line_spacing=1.35)


def slide_22_other_escalations():
    s = new_slide(22)
    add_eyebrow(s, "21  ·  other escalations")
    add_headline(s, "Three more situations trigger an owner page.", size=32)
    cards = [
        ("Silent query",
         "Customer asks something Sunny can't safely answer (exact install date, unusual brand, warranty claim). Owner gets the question. Sunny gives the customer a holding line."),
        ("Negotiation",
         "Customer pushes on price past Sunny's 'this is our best price' line. Owner decides whether to discount."),
        ("Repeat or complex",
         "Same customer asks the same thing again, or asks for something multi-step. Owner takes over manually."),
    ]
    for i, (title, body) in enumerate(cards):
        x = 0.7 + i * 4.1
        add_round_rect(s, x, 3.1, 3.85, 3.7, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
        add_text(s, title, x + 0.35, 3.4, 3.4, 0.7, size=22, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body, x + 0.35, 4.15, 3.4, 2.5, size=14, color=INK_SOFT, line_spacing=1.45)


def slide_23_alert_anatomy():
    s = new_slide(23)
    add_eyebrow(s, "22  ·  anatomy of an alert")
    add_headline(s, "Every alert carries enough context to act in ten seconds.", size=28, h=1.1)
    # left: real admin owner-chat screenshot
    add_image_card(s, ADMIN_OWNER_CHAT_ONE, 0.7, 2.45, 7.4, 4.0, label="OWNER CHAT  ·  REAL HOT ALERT")
    # annotations on right
    add_text(s, "WHAT'S INSIDE EVERY ALERT", 8.4, 2.55, 4.2, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    items = [
        ("01", "HOT LEAD header — color-coded by alert type."),
        ("02", "Customer name and phone."),
        ("03", "Lead Category, Temp, Intent (pickup, payment, technical)."),
        ("04", "Latest message verbatim."),
        ("05", "Six-turn conversation summary."),
        ("06", "Admin deep-link to that conversation."),
        ("07", "Open WhatsApp click-to-chat link."),
    ]
    y = 2.95
    for i, (num, it) in enumerate(items):
        add_text(s, num, 8.4, y + i * 0.48, 0.4, 0.35, size=12, color=EMBER, bold=True)
        add_text(s, it, 8.85, y + i * 0.48, 3.75, 0.4, size=11, color=INK_SOFT, line_spacing=1.3)
    add_text(s,
             "This is a real alert from Owner Chat. Two HOT LEADs are visible: one for a customer 'Ade' (intent: pickup_question) and one for 'Obara' (intent: payment_question). The owner can act on either one from his phone without leaving WhatsApp.",
             0.7, 6.6, 11.933, 0.4, size=10, color=INK_MUTE, italic=True, line_spacing=1.35)


def slide_24_admin_overview():
    s = new_slide(24)
    add_eyebrow(s, "23  ·  admin overview")
    add_headline(s, "One web panel. Everything the owner needs.", size=32)
    tabs = ["Inbox", "Contacts", "Warehouse Stock", "Owner Chat", "Rules editor", "Models & Config", "Reports"]
    cell_w, cell_h = 1.62, 0.7
    gap = 0.1
    total_w = len(tabs) * cell_w + (len(tabs) - 1) * gap
    start_x = (13.333 - total_w) / 2
    y = 3.3
    for i, t in enumerate(tabs):
        x = start_x + i * (cell_w + gap)
        add_round_rect(s, x, y, cell_w, cell_h, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.1)
        add_text(s, t, x + 0.1, y + 0.1, cell_w - 0.2, cell_h - 0.2, size=12, color=INK, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    add_text(s, "WhatsApp-style two-pane UI. Light theme. Mobile-friendly. Edits hit Sunny on the next customer message. No redeploy.",
             0.7, 4.6, 11.933, 1.5, size=18, color=INK_SOFT, line_spacing=1.45, align=PP_ALIGN.CENTER)


def admin_real_slide(page, eyebrow_text, headline, image_path, image_label,
                     explanation, key_callouts=None, headline_size=28):
    """Admin tab slide with real screenshot and an explanation paragraph.

    explanation: a paragraph describing what's IN the screenshot.
    key_callouts: optional list of short bullets highlighting specific UI elements.
    """
    s = new_slide(page)
    add_eyebrow(s, eyebrow_text)
    add_headline(s, headline, size=headline_size, h=1.2)
    # left: real screenshot card
    add_image_card(s, image_path, 0.7, 2.4, 7.7, 4.3, label=image_label)
    # right: explanation panel
    add_text(s, "WHAT YOU'RE LOOKING AT", 8.7, 2.4, 4.0, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s, explanation, 8.7, 2.8, 3.95, 2.2, size=12, color=INK_SOFT, line_spacing=1.45)
    if key_callouts:
        add_text(s, "KEY ELEMENTS", 8.7, 5.05, 4.0, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
        y = 5.4
        for b in key_callouts:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.7), Inches(y + 0.1), Inches(0.08), Inches(0.08))
            dot.fill.solid(); dot.fill.fore_color.rgb = SUN; dot.line.fill.background(); dot.shadow.inherit = False
            add_text(s, b, 8.9, y, 3.75, 0.32, size=11, color=INK, line_spacing=1.3)
            y += 0.42


def admin_mock_slide(page, eyebrow_text, headline, mock_title, mock_lines, side_bullets, headline_size=28):
    s = new_slide(page)
    add_eyebrow(s, eyebrow_text)
    add_headline(s, headline, size=headline_size, h=1.2)
    # left: placeholder
    placeholder_mock(s, 0.7, 2.7, 7.5, 4.05, mock_title, mock_lines, mock_label="MOCK")
    # right: bullets
    add_text(s, "WHAT YOU SEE", 8.5, 2.7, 4.1, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    y = 3.05
    for b in side_bullets:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(y + 0.13), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = SUN; dot.line.fill.background(); dot.shadow.inherit = False
        add_text(s, b, 8.75, y, 3.85, 0.5, size=13, color=INK_SOFT, line_spacing=1.35)
        y += 0.55


def add_image_card(slide, image_path, x, y, w, h, label=None):
    """Embed a real image inside a clean card frame. Auto-fits while preserving aspect ratio.

    Args:
      x, y, w, h: card bounding box in inches.
      label: optional small uppercase label above the image (e.g. 'CUSTOMER THREAD').
    """
    # Outer card
    add_round_rect(slide, x, y, w, h, WHITE, line=LINE_STRONG, line_w=0.75, radius_in=0.15)
    inset = 0.18
    label_h = 0.32 if label else 0
    img_x = x + inset
    img_y = y + inset + label_h
    img_w_max = w - 2 * inset
    img_h_max = h - 2 * inset - label_h
    if label:
        add_text(slide, label, x + inset, y + inset, w - 2 * inset, 0.25,
                 size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    # Insert picture so it fits the bounding box, preserving aspect ratio
    # python-pptx will use image's native aspect if only one of w/h is given;
    # so we figure out aspect from the file and pick the limiting dimension.
    from PIL import Image as PILImage
    with PILImage.open(image_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    if img_w_max / img_h_max > aspect:
        # height is the limit
        ph = img_h_max
        pw = ph * aspect
        px = img_x + (img_w_max - pw) / 2
        py = img_y
    else:
        # width is the limit
        pw = img_w_max
        ph = pw / aspect
        px = img_x
        py = img_y + (img_h_max - ph) / 2
    slide.shapes.add_picture(image_path, Inches(px), Inches(py), Inches(pw), Inches(ph))


def placeholder_mock(slide, x, y, w, h, title, lines, mock_label="MOCK"):
    """Generic clean mock: outer card + title + line items. No layered overlapping UI."""
    add_round_rect(slide, x, y, w, h, WHITE, line=LINE_STRONG, line_w=0.75, radius_in=0.15)
    # [MOCK] badge top right
    bw = 0.65
    add_round_rect(slide, x + w - bw - 0.12, y + 0.12, bw, 0.3, EMBER, line=None, radius_in=0.15)
    add_text(slide, mock_label, x + w - bw - 0.12, y + 0.12, bw, 0.3, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE, letter_spacing_pts=1.5)
    # title
    add_text(slide, title, x + 0.25, y + 0.2, w - bw - 0.6, 0.35, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    # lines
    line_y = y + 0.7
    line_h = max(0.32, (h - 0.85) / max(1, len(lines)))
    for label, body in lines:
        if label:
            add_text(slide, label.upper(), x + 0.25, line_y, 1.6, 0.3, size=9, color=INK_MUTE, bold=True, letter_spacing_pts=1.5, va=MSO_ANCHOR.TOP)
            add_text(slide, body, x + 1.95, line_y, w - 2.2, line_h - 0.05, size=13, color=INK, line_spacing=1.35, va=MSO_ANCHOR.TOP)
        else:
            add_text(slide, body, x + 0.25, line_y, w - 0.5, line_h - 0.05, size=13, color=INK_SOFT, line_spacing=1.35, va=MSO_ANCHOR.TOP)
        line_y += line_h


def slide_25_inbox():
    admin_real_slide(25, "24  ·  inbox", "Every conversation, in one place.",
                     ADMIN_INBOX, "ADMIN  ·  INBOX TAB",
                     explanation=(
                         "Two-pane WhatsApp-style view. The left column lists every active conversation, "
                         "color-coded by category (HOT in red, SERIOUS in amber, COLD in grey). "
                         "Right pane shows the selected thread: in this snapshot, Michael (a SERIOUS installer) "
                         "is being advised on a Growatt-to-Deye inverter replacement with 10 kVA storage. "
                         "Sunny is replying autonomously. The operator can click Take Over to step in at any moment."
                     ),
                     key_callouts=[
                         "Top bar: 113 pending queries waiting",
                         "Avatars color-coded by category",
                         "Take Over button on every conversation",
                         "Reply box at the bottom for manual replies",
                     ])


def slide_26_contacts():
    admin_real_slide(26, "25  ·  contacts", "Every customer, sortable, exportable.",
                     ADMIN_CONTACTS, "ADMIN  ·  CONTACTS TAB",
                     explanation=(
                         "Over 1,023 contacts so far. Each row carries the classifier's verdict: category, "
                         "temperature, client type (residential, installer, integrator), language, and location. "
                         "All filled in automatically as customers message in. "
                         "Filters at the top slice by category, temperature, or active-window. "
                         "The green Export button (top-right) downloads every row as Excel in one click."
                     ),
                     key_callouts=[
                         "Live count of total contacts in the header",
                         "Category column (HOT, SERIOUS, COLD, DISQUALIFIED)",
                         "Phone column preserves leading zeros (text format)",
                         "Click any row to open that conversation",
                     ])


def slide_27_warehouse_tab():
    admin_real_slide(27, "26  ·  warehouse stock",
                     "Live inventory. The only place Sunny gets prices.",
                     ADMIN_WAREHOUSE, "ADMIN  ·  WAREHOUSE STOCK",
                     explanation=(
                         "The source of truth for everything Sunny says about stock and pricing. "
                         "Items are grouped by category (BATTERIES, CONTROL BOX, RACK). "
                         "Each row shows live stock state per location (Abuja and Lagos), unit count, and price in NGN. "
                         "Incoming items carry a delivery note. "
                         "The owner toggles state in real time; the next customer reply uses the new state, no redeploy."
                     ),
                     key_callouts=[
                         "Independent Abuja and Lagos stock per item",
                         "In stock, Incoming, or Out of stock",
                         "Per-item datasheet PDF attachable",
                         "Per-item photos (JPG / PNG)",
                     ])


def slide_28_owner_chat():
    admin_real_slide(28, "27  ·  owner chat",
                     "Every alert, every reply, in one read-only thread.",
                     ADMIN_OWNER_CHAT_MANY, "ADMIN  ·  OWNER CHAT",
                     explanation=(
                         "A full audit trail of every message Sunny ever sends to the owner. "
                         "Each HOT LEAD alert is preserved with customer name, phone, lead category, "
                         "intent (payment_question, technical_question, pickup_question), the latest message, "
                         "and a 6-turn conversation summary. The owner's own replies are interleaved so the context never disappears."
                     ),
                     key_callouts=[
                         "ESCALATION_ALERT_HOT header on each entry",
                         "Customer phone + lead Category + Intent",
                         "Full conversation summary inside the alert",
                         "Open WhatsApp Chat click-to-chat link",
                     ])


def slide_29_rules_editor():
    admin_real_slide(29, "28  ·  rules editor",
                     "Edit Sunny's brain in plain English. Save. Deploy.",
                     ADMIN_RULES, "ADMIN  ·  RULES EDITOR",
                     explanation=(
                         "Sunny's personality, sales doctrine, voice rules, and hard limits all live in three editable markdown files: "
                         "system.md (the master prompt, visible here), classifier.md (how leads are categorized), "
                         "and owner_qa.md (boss-mode answers). The owner types changes directly in the browser textarea, "
                         "hits Save, and the file commits to GitHub. Sunny re-reads his prompts every 30 seconds, "
                         "so changes apply on the next customer message with no restart."
                     ),
                     key_callouts=[
                         "Three live-editable prompt files",
                         "Save commits to GitHub automatically",
                         "Deploy to live triggers Railway redeploy",
                         "No developer involvement required",
                     ])


def slide_30_models_config():
    admin_real_slide(30, "29  ·  models & config",
                     "Live spending. One model. Total transparency.",
                     ADMIN_MODELS, "ADMIN  ·  MODELS & CONFIG",
                     explanation=(
                         "Real-time spending shown to the cent. Today, this month, daily budget cap. "
                         "All three call sites (Reply, Classifier, Owner teaching) run on the same model: Claude Sonnet 4.6. "
                         "Below: runtime configuration (DB path, media directory, daily budget) and WhatsApp environment (WABA ID, Graph API version). "
                         "If costs ever spike, the owner can flip any call site to a cheaper variant from this panel without touching code."
                     ),
                     key_callouts=[
                         "Today $7.18  ·  Month $142.41  ·  Budget $50",
                         "Reply: claude-sonnet-4-6",
                         "Classifier: claude-sonnet-4-6",
                         "Owner teaching: claude-sonnet-4-6",
                     ])


def slide_31_takeover():
    s = new_slide(31)
    add_eyebrow(s, "30  ·  human in the loop")
    add_headline(s, "The owner can jump into any conversation. Sunny stops typing.", size=28, h=1.1)
    steps = [
        ("01", "Take over", "Owner clicks 'Take over' on a conversation in the inbox."),
        ("02", "Sunny pauses", "Sunny stops auto-replying on that conversation only."),
        ("03", "Owner replies manually", "Owner types replies as if from Sunny. Customer sees them as normal."),
        ("04", "Release", "Owner clicks 'Release', or walks away (auto-release kicks in)."),
    ]
    y = 3.0
    for i, (num, title, body) in enumerate(steps):
        x = 0.7 + i * 3.1
        add_round_rect(s, x, y, 2.9, 3.5, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
        add_text(s, num, x + 0.25, y + 0.25, 1.0, 0.5, size=24, color=EMBER, bold=True, font=HEAD_FONT)
        add_text(s, title, x + 0.25, y + 0.8, 2.5, 0.6, size=18, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body, x + 0.25, y + 1.5, 2.5, 1.8, size=13, color=INK_SOFT, line_spacing=1.45)
        if i < 3:
            add_arrow(s, x + 2.95, y + 1.75, x + 3.06, y + 1.75, color=INK_SOFT, weight=1.5)


def slide_32_autorelease():
    s = new_slide(32)
    add_eyebrow(s, "31  ·  auto-release")
    add_headline(s, "Sunny resumes when the owner steps away.", size=32)
    boxes = [
        ("Human handles",        "Owner takes over a thread."),
        ("15 minutes idle",      "No new manual reply, auto-release cron checks."),
        ("Auto-release",         "Conversation flips back to Sunny."),
        ("Re-queue + reply",     "Last unanswered customer message is re-classified and answered."),
    ]
    y = 3.3
    for i, (title, body) in enumerate(boxes):
        x = 0.7 + i * 3.1
        add_round_rect(s, x, y, 2.9, 2.5, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
        add_text(s, title, x + 0.25, y + 0.3, 2.5, 0.5, size=17, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body, x + 0.25, y + 0.95, 2.5, 1.45, size=13, color=INK_SOFT, line_spacing=1.45)
        if i < 3:
            add_arrow(s, x + 2.95, y + 1.25, x + 3.06, y + 1.25, color=INK_SOFT, weight=1.5)
    add_text(s, "Tunable per deployment via HUMAN_AUTO_RELEASE_MINUTES env var. Default is 15 minutes.",
             0.7, 6.4, 11.933, 0.4, size=14, color=INK_MUTE, italic=True, align=PP_ALIGN.CENTER)


def slide_33_technical_reply():
    s = new_slide(33)
    add_eyebrow(s, "32  ·  technical reply")
    add_headline(s, "Sizing, address, datasheet. All from one customer message.", size=26, h=1.1)
    add_image_card(s, WA_TECHNICAL, 4.5, 2.3, 4.3, 4.55, label="REAL WHATSAPP THREAD")
    # explanation panel on left
    add_text(s, "WHAT SUNNY DID", 0.7, 2.4, 3.6, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s,
             "Shared the Lagos office address from his locations data, "
             "verbatim and accurate.",
             0.7, 2.85, 3.6, 0.8, size=13, color=INK_SOFT, line_spacing=1.4)
    add_text(s,
             "Refused to invent specs when the customer asked vaguely "
             "about 'a 6kw inverter'. Asked one clarifying question.",
             0.7, 3.75, 3.6, 1.1, size=13, color=INK_SOFT, line_spacing=1.4)
    add_text(s,
             "Once the customer clarified, quoted the exact model from "
             "warehouse: Deye SUN-6K-OG01LP1-EU-AM2. Topology and stock.",
             0.7, 4.95, 3.6, 1.2, size=13, color=INK_SOFT, line_spacing=1.4)
    add_text(s,
             "Sent the actual datasheet PDF as a WhatsApp document when "
             "asked.",
             0.7, 6.25, 3.6, 0.6, size=13, color=INK_SOFT, line_spacing=1.4)
    # right side caption
    add_text(s,
             "No invented model. No invented price. Sourced from warehouse on every reply.",
             9.0, 2.85, 3.6, 4.0, size=14, color=INK, italic=True, line_spacing=1.5)


def slide_34_costing_reply():
    s = new_slide(34)
    add_eyebrow(s, "33  ·  stock awareness")
    add_headline(s, "Three back-to-back stock questions. Three honest answers.", size=26, h=1.1)
    add_image_card(s, WA_STOCK, 4.5, 2.3, 4.3, 4.55, label="REAL WHATSAPP THREAD")
    # left annotations
    add_text(s, "WHAT THE CUSTOMER ASKED", 0.7, 2.4, 3.6, 0.3, size=10, color=EMBER, bold=True, letter_spacing_pts=2.0)
    questions = [
        ("01", "Do you have the 50kW HV?",
         "Sunny: Yes, Deye SUN-50K-SG01HP3-EU-BM4 (50kW HV, 3-phase) is available."),
        ("02", "Is the 16kWh in stock?",
         "Sunny: Deye SE-F16 is incoming, expected end of this week. You can pre-order to secure a unit."),
        ("03", "What about BOS-G?",
         "Sunny: Deye BOS-G-PACK 5.1 is out of stock right now."),
    ]
    y = 2.85
    for num, q, a in questions:
        add_text(s, num, 0.7, y, 0.4, 0.3, size=14, color=EMBER, bold=True)
        add_text(s, q, 1.1, y, 3.2, 0.3, size=13, color=INK, bold=True)
        add_text(s, a, 1.1, y + 0.32, 3.2, 0.9, size=11, color=INK_SOFT, line_spacing=1.35)
        y += 1.35
    # right caption
    add_text(s,
             "Honesty by default. No fake availability. No invented ETAs. The customer trusts what Sunny says because warehouse is the only source.",
             9.0, 2.85, 3.6, 4.0, size=13, color=INK, italic=True, line_spacing=1.5)


def slide_35_classification():
    s = new_slide(35)
    add_eyebrow(s, "34  ·  classification in action")
    add_headline(s, "'I want to pay'  →  classifier flips to HOT  →  owner is paged.", size=22, h=1.0)
    # left: customer thread
    add_image_card(s, WA_PANELS_HOT, 0.7, 2.35, 3.6, 4.05, label="CUSTOMER  ·  WHATSAPP")
    # right: owner alert
    add_image_card(s, ADMIN_OWNER_CHAT_ONE, 4.5, 2.35, 8.13, 4.05, label="OWNER  ·  ADMIN  ·  ALERT")
    add_text(s,
             "Left: customer browses solar panels, eventually says 'I want to pay'. Sunny acknowledges and hands off to the specialist with a wa.me link.   "
             "Right: at the same moment, the classifier upgrades this conversation to HOT (intent: payment_question) and the owner's alert lands in Owner Chat.   "
             "Two systems, one orchestrated handoff, end-to-end in under ten seconds.",
             0.7, 6.55, 11.933, 0.5, size=10, color=INK_MUTE, italic=True, line_spacing=1.35)


def slide_36_multilingual():
    s = new_slide(36)
    add_eyebrow(s, "35  ·  multilingual")
    add_headline(s, "Voice notes in Hausa, small talk in Pidgin, business in English.", size=28, h=1.1)
    exchanges = [
        ("English", "Customer: How much is the Deye 16kW?", "Sunny: ₦4,200,000. In stock in Lagos."),
        ("Pidgin",  "Customer: Bros, dat one don too much o.", "Sunny: I hear you. Make I check with my oga and revert."),
        ("Hausa",   "Customer: (voice note) Ina son inverter 5 kVA.", "Sunny: Mun da Deye 5kVA. Naira 1,250,000. Akwai a Abuja."),
    ]
    y = 2.9
    for i, (lang, cust, sun) in enumerate(exchanges):
        ry = y + i * 1.3
        add_round_rect(s, 0.7, ry, 11.933, 1.15, PAPER, line=LINE_STRONG, line_w=0.4, radius_in=0.1)
        add_text(s, lang.upper(), 0.95, ry + 0.12, 1.5, 0.22, size=11, color=EMBER, bold=True, letter_spacing_pts=2.0)
        add_text(s, cust, 0.95, ry + 0.42, 11.5, 0.3, size=13, color=INK, va=MSO_ANCHOR.TOP)
        add_text(s, sun,  0.95, ry + 0.78, 11.5, 0.3, size=13, color=INK_SOFT, italic=True, va=MSO_ANCHOR.TOP)


def slide_37_owner_view():
    s = new_slide(37)
    add_eyebrow(s, "36  ·  owner view")
    add_headline(s, "Here is what hits the owner's phone.", size=36, h=1.0)
    add_image_card(s, ADMIN_OWNER_CHAT_MANY, 0.7, 2.35, 12.0, 4.05, label="OWNER CHAT  ·  AS THE BOSS SEES IT")
    add_text(s,
             "Multiple HOT LEAD alerts in sequence. For each one, the owner sees: who, what they want, the lead category and temperature, the intent (payment, technical, pickup), the customer's last message, a 6-turn conversation summary, and one-click links to admin and to the customer's WhatsApp.",
             0.7, 6.55, 11.933, 0.4, size=10, color=INK_MUTE, italic=True, line_spacing=1.35, align=PP_ALIGN.CENTER)


def slide_38_comparison():
    s = new_slide(38)
    add_eyebrow(s, "37  ·  the difference, summarized")
    add_headline(s, "Side by side.", size=36, h=0.9)
    rows = [
        ("Memory",        "Stateless",         "Stateful, structured"),
        ("Knowledge",     "Static FAQ",        "Live warehouse + facts"),
        ("Languages",     "One",               "Five (including Pidgin)"),
        ("Voice notes",   "Ignored",           "Transcribed + handled"),
        ("Photos",        "Not handled",       "Vision + sent from warehouse"),
        ("Documents",     "Not sent",          "Datasheets auto-sent"),
        ("Escalation",    "None",              "HOT, silent, negotiation, repeat"),
        ("Self-check",    "None",              "14 reply guards per reply"),
        ("Owner control", "None",              "Live admin, prompts, models"),
    ]
    y0 = 2.4
    # Header
    add_text(s, "DIMENSION", 0.7, y0, 3.0, 0.3, size=10, color=INK_MUTE, bold=True, letter_spacing_pts=2.0)
    add_text(s, "CHATBOT",   4.2, y0, 4.0, 0.3, size=10, color=INK_MUTE, bold=True, letter_spacing_pts=2.0)
    add_text(s, "SUNNY",     8.6, y0, 4.0, 0.3, size=10, color=EMBER,    bold=True, letter_spacing_pts=2.0)
    add_line(s, 0.7, y0 + 0.32, 12.633, y0 + 0.32, color=LINE_STRONG, weight=0.75)
    for i, (dim, cb, sn) in enumerate(rows):
        ry = y0 + 0.45 + i * 0.42
        add_text(s, dim, 0.7, ry, 3.0, 0.4, size=14, color=INK, bold=True, font=HEAD_FONT, va=MSO_ANCHOR.TOP)
        add_text(s, cb,  4.2, ry, 4.0, 0.4, size=13, color=INK_MUTE, va=MSO_ANCHOR.TOP)
        add_text(s, sn,  8.6, ry, 4.0, 0.4, size=13, color=INK, va=MSO_ANCHOR.TOP)
        add_line(s, 0.7, ry + 0.36, 12.633, ry + 0.36, color=LINE, weight=0.3)


def slide_39_pillars():
    s = new_slide(39)
    add_eyebrow(s, "38  ·  the pillars")
    add_headline(s, "Seven pillars that make Sunny different.", size=32)
    pillars = [
        ("01", "Rule discipline",   "14 reply guards. No exceptions."),
        ("02", "Stateful memory",   "Conversation state engine on every turn."),
        ("03", "Multi-language",    "Five languages, first-class."),
        ("04", "Multimodal",        "Text, voice, photos, documents."),
        ("05", "Classification",    "Every lead, categorized in real time."),
        ("06", "Escalation",        "Owner paged with full context."),
        ("07", "Admin control",     "Live edits. No developer loop."),
    ]
    # 4 + 3 grid
    cell_w, cell_h = 2.95, 1.65
    gap = 0.15
    start_x = 0.7
    y = 3.0
    for i, (num, title, body) in enumerate(pillars):
        if i < 4:
            x = start_x + i * (cell_w + gap)
            ry = y
        else:
            x = start_x + (i - 4) * (cell_w + gap) + (cell_w + gap) / 2
            ry = y + cell_h + gap
        add_round_rect(s, x, ry, cell_w, cell_h, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.12)
        add_text(s, num, x + 0.25, ry + 0.2, 1.0, 0.3, size=14, color=EMBER, bold=True, letter_spacing_pts=2.0)
        add_text(s, title, x + 0.25, ry + 0.6, cell_w - 0.5, 0.45, size=17, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body, x + 0.25, ry + 1.1, cell_w - 0.5, 0.5, size=12, color=INK_SOFT, line_spacing=1.35)


def slide_40_cost():
    s = new_slide(40)
    add_eyebrow(s, "39  ·  cost economics")
    add_headline(s, "What Sunny costs ElectroSun. Actual numbers.", size=34)
    # two columns
    add_round_rect(s, 0.7, 3.0, 5.95, 3.6, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
    add_text(s, "SUNNY",       0.95, 3.2, 5.4, 0.3, size=11, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s, "$140 / month",   0.95, 3.5, 5.4, 0.9, size=48, color=INK, bold=True, font=HEAD_FONT)
    add_text(s, "Around $5 per day  ·  5,800+ replies / month", 0.95, 4.5, 5.4, 0.5, size=18, color=INK_SOFT)
    add_text(s, "Sonnet 4.6 across the whole pipeline. Daily budget of $50 with auto-fallback if exceeded. Prompt caching keeps the cost flat.",
             0.95, 5.1, 5.4, 1.4, size=14, color=INK_SOFT, line_spacing=1.4)

    add_round_rect(s, 6.85, 3.0, 5.78, 3.6, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
    add_text(s, "HUMAN EQUIVALENT", 7.1, 3.2, 5.4, 0.3, size=11, color=EMBER, bold=True, letter_spacing_pts=2.0)
    add_text(s, "₦600k / month",    7.1, 3.5, 5.4, 0.9, size=48, color=INK, bold=True, font=HEAD_FONT)
    add_text(s, "2 staff for 24/7 coverage",  7.1, 4.5, 5.4, 0.5, size=18, color=INK_SOFT)
    add_text(s, "One WhatsApp account manager covers maybe 12 hours a day. Two are needed for 24/7. Plus shift handoff overhead and supervision.",
             7.1, 5.1, 5.4, 1.4, size=14, color=INK_SOFT, line_spacing=1.4)

    add_text(s, "Sunny replaces about two full-time staff at one-tenth the monthly cost. Pays for himself in week one.",
             0.7, 6.75, 11.933, 0.3, size=12, color=INK_MUTE, italic=True, align=PP_ALIGN.CENTER)


def slide_41_whitelabel():
    s = new_slide(41)
    add_eyebrow(s, "40  ·  what's next")
    add_headline(s, "Same architecture. Any vertical.", size=40)
    # center box
    add_round_rect(s, 5.4, 2.95, 2.5, 1.0, INK, line=None, radius_in=0.15)
    add_text(s, "sunny-template", 5.4, 3.05, 2.5, 0.4, size=16, color=SUN_SOFT, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    add_text(s, "the white-label",  5.4, 3.45, 2.5, 0.4, size=11, color=CREAM, italic=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE, letter_spacing_pts=1.5)
    # three industry boxes
    industries = [
        ("Dental clinic", "Appointment booking,\nreminders, follow-ups."),
        ("Real estate",    "Listing inquiries,\nviewing scheduling."),
        ("Retail",         "Order status,\nphotos, returns."),
    ]
    iy = 4.7
    for i, (name, body) in enumerate(industries):
        x = 1.0 + i * 4.1
        add_round_rect(s, x, iy, 3.6, 1.5, PAPER, line=LINE_STRONG, line_w=0.5, radius_in=0.15)
        add_text(s, name, x + 0.3, iy + 0.15, 3.0, 0.4, size=18, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, body, x + 0.3, iy + 0.65, 3.0, 0.8, size=12, color=INK_SOFT, line_spacing=1.35)
        # arrow from center down to box
        add_arrow(s, 6.65, 4.0, x + 1.8, iy - 0.02, color=INK_SOFT, weight=1.0)

    add_text(s, "Eighty percent of Sunny is industry-agnostic plumbing. Twenty percent is the vertical pack. Each new client takes days, not months.",
             0.7, 6.55, 11.933, 0.3, size=12, color=INK_MUTE, italic=True, align=PP_ALIGN.CENTER)


def slide_42_cta():
    s = new_slide(42, show_chrome=False)
    add_rect(s, 0, 0, 13.333, 7.5, CREAM, line=None)
    # decorative sun in top-right, fully on canvas
    sun_big = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(0.4), Inches(2.2), Inches(2.2))
    sun_big.fill.solid(); sun_big.fill.fore_color.rgb = SUN_SOFT; sun_big.line.fill.background(); sun_big.shadow.inherit = False

    add_text(s, "BUILD YOURS", 0.7, 1.8, 6, 0.35, size=12, color=EMBER, bold=True, letter_spacing_pts=4.0)
    add_text(s, "Want a Sunny for", 0.7, 2.4, 11, 1.0, size=60, color=INK, bold=True, font=HEAD_FONT, line_spacing=1.05)
    add_text(s, "your business?",   0.7, 3.55, 11, 1.0, size=60, color=INK, bold=True, italic=True, font=HEAD_FONT, line_spacing=1.05)
    add_text(s,
             "We deploy custom WhatsApp account managers for businesses that take inbound on WhatsApp. "
             "Solar, retail, services, real estate, anything. From kickoff to live in two weeks.",
             0.7, 5.1, 10.5, 1.4, size=20, color=INK_SOFT, line_spacing=1.45)

    add_line(s, 0.7, 6.8, 12.633, 6.8, color=LINE_STRONG, weight=0.5)
    add_text(s, "[ your email ]   ·   [ your website ]   ·   [ your WhatsApp ]",
             0.7, 6.95, 11.933, 0.35, size=14, color=INK_MUTE, align=PP_ALIGN.CENTER)


# -----------------------------------------------------------------------------
# Build all slides
# -----------------------------------------------------------------------------

builders = [
    slide_01_cover, slide_02_pitch, slide_03_opportunity, slide_04_before, slide_05_chatbot_fails,
    slide_06_capabilities, slide_07_not_chatbot, slide_08_day, slide_09_numbers,
    slide_10_architecture, slide_11_where_lives, slide_12_brain, slide_13_pipeline, slide_14_state, slide_15_guards,
    slide_16_warehouse, slide_17_knowledge, slide_18_owner_qa, slide_19_multimodal,
    slide_20_lead_cats, slide_21_hot_handoff, slide_22_other_escalations, slide_23_alert_anatomy,
    slide_24_admin_overview,
    slide_25_inbox, slide_26_contacts, slide_27_warehouse_tab, slide_28_owner_chat, slide_29_rules_editor, slide_30_models_config,
    slide_31_takeover, slide_32_autorelease,
    slide_33_technical_reply, slide_34_costing_reply, slide_35_classification, slide_36_multilingual, slide_37_owner_view,
    slide_38_comparison, slide_39_pillars, slide_40_cost,
    slide_41_whitelabel, slide_42_cta,
]

for fn in builders:
    fn()

OUTPUT = Path(__file__).resolve().parents[1] / "presentation" / "sunny-case-study.pptx"
prs.save(OUTPUT)
print(f"wrote {OUTPUT}  ·  {len(prs.slides)} slides")
